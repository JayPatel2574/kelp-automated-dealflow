from pptx.util import Inches
from pptx.enum.text import PP_ALIGN

from src.ppt.theme import FONT_BODY, FOOTER_SIZE, LOGO_SIZE
from src.ppt.colors import KELP_DARK_GREY


FOOTER_TEXT = "Strictly Private & Confidential – Prepared by Kelp M&A Team"
LOGO_TEXT = "Kelp"


def add_footer(slide):
    box = slide.shapes.add_textbox(
        Inches(1.5), Inches(6.9),
        Inches(7), Inches(0.3)
    )
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = FOOTER_TEXT
    p.font.name = FONT_BODY
    p.font.size = FOOTER_SIZE
    p.font.color.rgb = KELP_DARK_GREY
    p.alignment = PP_ALIGN.CENTER


def add_logo_text(slide):
    box = slide.shapes.add_textbox(
        Inches(0.3), Inches(0.2),
        Inches(2), Inches(0.4)
    )
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = LOGO_TEXT
    p.font.name = FONT_BODY
    p.font.size = LOGO_SIZE
    p.font.bold = True
    p.font.color.rgb = KELP_DARK_GREY
    p.alignment = PP_ALIGN.LEFT
