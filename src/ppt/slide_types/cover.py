from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from src.ppt.colors import KELP_DARK_INDIGO, KELP_WHITE
from src.ppt.branding import add_logo_text, add_footer
from src.ppt.theme import FONT_HEADING, TITLE_SIZE
from src.ppt.layout_engine import full_bleed_background


def build_cover(prs, slide, sector):
    # Dark indigo full bleed
    full_bleed_background(slide, prs, KELP_DARK_INDIGO)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(1.2), Inches(2.5),
        Inches(7), Inches(1.5)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Investment Brief"
    p.font.name = FONT_HEADING
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = KELP_WHITE
    p.alignment = PP_ALIGN.LEFT

    # Subtitle
    p2 = tf.add_paragraph()
    p2.text = f"Sector: {sector.replace('_', ' ')}"
    p2.font.size = Pt(14)
    p2.font.color.rgb = KELP_WHITE

    add_logo_text(slide)
    add_footer(slide)
