from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE
from src.ppt.colors import KELP_DARK_INDIGO, KELP_WHITE


def full_bleed_background(slide, prs, color=KELP_DARK_INDIGO):
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0,
        prs.slide_width,
        prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()


def white_content_panel(slide, left, top, width, height):
    panel = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left, top,
        width, height
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = KELP_WHITE
    panel.line.fill.background()
    return panel
