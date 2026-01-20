from pptx import Presentation

from src.ppt.slide_types.cover import build_cover
from src.ppt.slide_types.overview import build_overview
from src.ppt.slide_types.highlights import build_highlights
from src.ppt.slide_types.financials import build_financials
from src.ppt.slide_types.disclaimer import build_disclaimer
from src.ai.investment_narrator import InvestmentNarrator

class PPTBuilder:

    def __init__(self):
        self.prs = Presentation()

    def build(self, data, insights, sector, output_path):
        # Cover
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        build_cover(self.prs, slide, sector)

        narrator = InvestmentNarrator()

        # Overview
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        overview_text = narrator.generate_overview(data, sector)
        build_overview(self.prs, slide, data, overview_text)

        # Investment Highlights
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        highlights_ai = narrator.generate_highlights(data, sector)
        build_highlights(self.prs, slide, highlights_ai)

        # Financial Summary
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        build_financials(self.prs, slide, data)

        # Disclaimer
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        build_disclaimer(self.prs, slide)

        self.prs.save(output_path)
